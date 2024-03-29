
# from PyPDF2 import PdfReader
import PyPDF2
import docx
import pptx
import csv
from openai import OpenAI
import base64
import textwrap

from pytube import YouTube 
import whisper

def get_plain_text_from_txt(txt_filename):
    # Initialize empty string 
    txt_text = ""

    # Open text file
    with open(txt_filename, 'r') as f:
        # Loop through lines
        for line in f:
            # Append line to the string
            txt_text += line

    # Return text
    return txt_text

def get_plain_text(pdf_filename):
    # Open the PDF file in read-binary mode
    with open(pdf_filename, 'rb') as f:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(f)
        # Get the number of pages in the PDF file
        num_pages =  len(pdf_reader.pages)
        # Initialize an empty string to store the text
        document_text = ''
        # Loop through each page in the PDF file
        for page_num in range(num_pages):
        # Get the current page
            page = pdf_reader.pages[page_num]
            #Extract the text from the current page
            page_text = page.extract_text()
            # Append the page text to the document text
            document_text += page_text
        return document_text

def get_plain_text_docx(docx_filename):

  # Open the Word document
  document = docx.Document(docx_filename)  

  # Initialize empty string
  doc_text = ""

  # Loop through paragraphs
  for para in document.paragraphs:
    doc_text += para.text + "\n"  
  
  # Loop through tables
  for table in document.tables:
    for row in table.rows:
      for cell in row.cells:
        doc_text += cell.text + "\n"

  # Return the text
  return doc_text




def  get_plain_text_pptx(pptx_filename):

  # Open the PowerPoint presentation
  presentation = pptx.Presentation(pptx_filename)

  # Initialize empty string 
  text = ""

  # Loop through slides 
  for slide in presentation.slides:

    # Extract text from each paragraph of the slide
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text += shape.text + "\n"

  # Return the text
  return text  



def get_plain_text_csv(csv_filename):

  # Open CSV file
  with open(csv_filename, 'r') as f:

    # Create CSV reader
    reader = csv.reader(f)

    # Skip header row
    # next(reader)  

    # Initialize empty string 
    csv_text = ""

    # Loop through rows
    for row in reader:
    
      # Append row values 
      csv_text += " ".join(row) + "\n"

  # Return text
  return csv_text

def get_plain_text_mp4(mp4_url):
    
    try: 
        # object creation using url to mp4 location
        yt = YouTube(mp4_url) 
    except: 
        print("Connection Error")

    # Download video stream to file
    
    yt.streams.filter(file_extension='mp4')
    stream = yt.streams.get_by_itag(139)
    stream.download('',"AudioTranscript.mp4")

    model = whisper.load_model("base.en")
    audio_file= "AudioTranscript.mp4"
    result = model.transcribe(audio_file)

    result_text = result['text']
    

    wrapped_text = textwrap.fill(result_text, width=80)
    return wrapped_text
    
def get_plain_text_mp3(mp3_url):
   
    model = whisper.load_model("base.en")
    audio_file= mp3_url
    result = model.transcribe(audio_file)

    result_text = result['text']

    wrapped_text = textwrap.fill(result_text, width=80)
    return wrapped_text


def split_text_into_chunks(plain_text, max_chars=2000):
    text_chunks = []
    current_chunk = ""
    for line in plain_text.split("\n"):
        if len(current_chunk) + len(line) + 1 <= max_chars:
            current_chunk += line + " "
        else:
            text_chunks.append(current_chunk.strip())
            current_chunk = line + " "
    if current_chunk:
        text_chunks.append(current_chunk.strip())
    return text_chunks


# Function to encode the image
def encode_image(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def get_plain_text_image(image_file):
    client = OpenAI()

    # Path to your image
    image_path = image_file   #"c:/temp/data-model-01.png"

    # Getting the base64 string
    base64_image = encode_image(image_path)

    response = client.chat.completions.create(
    model="gpt-4-vision-preview",
    messages=[
        {
            "role": "user",
            "content": [
                {"type": "text", "text": "Analyse this image in detail, extract data in tabular form, label each table and provide a narrative for each table"},
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{base64_image}",
                    },
                },
            ],
        }
    ],
    max_tokens=1000,
    )
    return response.choices[0].message.content

def scrape_text_from_txt(txt_file, max_chars=2000):
    plain_text = get_plain_text_from_txt(txt_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_pdf(pdf_file, max_chars=2000):
    plain_text = get_plain_text(pdf_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_docx(docx_file, max_chars=2000):
    plain_text = get_plain_text_docx(docx_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_pptx(pptx_file, max_chars=2000):
    plain_text = get_plain_text_pptx(pptx_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_csv(csv_file, max_chars=2000):
    plain_text = get_plain_text_csv(csv_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks
 
def scrape_text_from_image(image_file, max_chars=2000):
    plain_text = get_plain_text_image(image_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_mp4(mp4_file, max_chars=2000):
    plain_text = get_plain_text_mp4(mp4_file)
    text_chunks = split_text_into_chunks(plain_text, max_chars)
    return text_chunks

def scrape_text_from_mp3(mp3_file, max_chars=2000):
    plain_text = get_plain_text_mp3(mp3_file)
    # text_chunks = split_text_into_chunks(plain_text, max_chars)
    # return text_chunks
    return plain_text
